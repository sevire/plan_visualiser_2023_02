"""
Create a sample PowerPoint with a realistic-looking visual timeline.
This simulates what the actual renderer would produce from a real visual.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class MockColor:
    def __init__(self, r, g, b):
        self.red = r
        self.green = g
        self.blue = b


class MockFont:
    def __init__(self, name="Arial"):
        self.font_name = name


class MockFormat:
    def __init__(self, fill_color, line_color, font_color, line_thickness=1, font_size=10):
        self.fill_color = fill_color
        self.line_color = line_color
        self.line_thickness = line_thickness
        self.font_color = font_color
        self.font_size = font_size
        self.font = MockFont()


class MockPlotable:
    """Mock of RectangleBasedPlotable"""
    def __init__(self, plotable_id, left, top, width, height, text="", shape_type='RECTANGLE',
                 fill_color=None, line_color=None, font_color=None, font_size=10):
        self.plotable_id = plotable_id
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.text = text

        # Default colors
        if fill_color is None:
            fill_color = MockColor(200, 220, 240)
        if line_color is None:
            line_color = MockColor(0, 0, 0)
        if font_color is None:
            font_color = MockColor(0, 0, 0)

        self.format = MockFormat(fill_color, line_color, font_color, font_size=font_size)

        class Shape:
            def __init__(self, shape_name):
                self.value = shape_name
        self.shape = Shape(shape_type)

        class TextFlow:
            FLOW_CENTRE = 'FLOW_CENTRE'
            FLOW_TO_LEFT = 'FLOW_TO_LEFT'
            FLOW_TO_RIGHT = 'FLOW_TO_RIGHT'

        class VerticalAlignment:
            TOP = 'TOP'
            MIDDLE = 'MIDDLE'
            BOTTOM = 'BOTTOM'

        self.text_flow = TextFlow.FLOW_CENTRE
        self.text_vertical_alignment = VerticalAlignment.MIDDLE
        self.external_text_flag = False

    def get_left(self):
        return self.left

    def get_top(self):
        return self.top

    def get_right(self):
        return self.left + self.width

    def get_bottom(self):
        return self.top + self.height

    def get_width(self):
        return self.width

    def get_height(self):
        return self.height


def create_realistic_visual():
    """Create a realistic project timeline visual"""

    # Timeline colors (light gray background)
    timeline_fill = MockColor(240, 240, 240)
    timeline_line = MockColor(100, 100, 100)
    timeline_font = MockColor(50, 50, 50)

    # Swimlane colors (light blue)
    swimlane_fill = MockColor(220, 235, 250)
    swimlane_line = MockColor(70, 130, 180)
    swimlane_font = MockColor(0, 0, 0)

    # Activity colors (various)
    activity_blue = MockColor(100, 150, 200)
    activity_green = MockColor(100, 200, 150)
    activity_orange = MockColor(250, 180, 100)
    activity_line = MockColor(50, 50, 50)
    activity_font = MockColor(0, 0, 0)

    # Create month timeline labels (Jan-Jun, 7 months)
    month_width = 150
    month_height = 25
    month_names = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul"]

    timelines = []
    for i, month in enumerate(month_names):
        timelines.append(MockPlotable(
            f"month_{i}",
            i * month_width, 0, month_width, month_height,
            month,
            fill_color=timeline_fill,
            line_color=timeline_line,
            font_color=timeline_font,
            font_size=12
        ))

    # Create quarter timeline labels
    quarter_height = 20
    timelines.append(MockPlotable(
        "q1",
        0, month_height, month_width * 3, quarter_height,
        "Q1 2024",
        fill_color=MockColor(200, 200, 200),
        line_color=timeline_line,
        font_color=timeline_font,
        font_size=11
    ))
    timelines.append(MockPlotable(
        "q2",
        month_width * 3, month_height, month_width * 4, quarter_height,
        "Q2 2024",
        fill_color=MockColor(220, 220, 220),
        line_color=timeline_line,
        font_color=timeline_font,
        font_size=11
    ))

    # Create swimlanes
    swimlane_top = month_height + quarter_height + 10
    swimlane_height = 80
    swimlane_gap = 10

    swimlanes = [
        MockPlotable(
            "swimlane1",
            0, swimlane_top, month_width * 7, swimlane_height,
            "Development Team",
            fill_color=swimlane_fill,
            line_color=swimlane_line,
            font_color=swimlane_font,
            font_size=11
        ),
        MockPlotable(
            "swimlane2",
            0, swimlane_top + swimlane_height + swimlane_gap, month_width * 7, swimlane_height,
            "QA Team",
            fill_color=MockColor(240, 220, 235),
            line_color=swimlane_line,
            font_color=swimlane_font,
            font_size=11
        ),
        MockPlotable(
            "swimlane3",
            0, swimlane_top + 2 * (swimlane_height + swimlane_gap), month_width * 7, swimlane_height,
            "Infrastructure",
            fill_color=MockColor(220, 250, 235),
            line_color=swimlane_line,
            font_color=swimlane_font,
            font_size=11
        )
    ]

    # Create activities
    activity_height = 35
    activities = []

    # Development activities
    y_offset = swimlane_top + 20
    activities.append(MockPlotable(
        "act1",
        month_width * 0.5, y_offset, month_width * 1.5, activity_height,
        "API Design",
        shape_type='ROUNDED_RECTANGLE',
        fill_color=activity_blue,
        line_color=activity_line,
        font_color=activity_font,
        font_size=10
    ))

    activities.append(MockPlotable(
        "act2",
        month_width * 2, y_offset, month_width * 2, activity_height,
        "Backend Development",
        shape_type='ROUNDED_RECTANGLE',
        fill_color=activity_blue,
        line_color=activity_line,
        font_color=activity_font,
        font_size=10
    ))

    activities.append(MockPlotable(
        "act3",
        month_width * 4, y_offset, month_width * 2.5, activity_height,
        "Frontend Development",
        shape_type='ROUNDED_RECTANGLE',
        fill_color=activity_green,
        line_color=activity_line,
        font_color=activity_font,
        font_size=10
    ))

    # Milestone
    milestone_size = 20
    activities.append(MockPlotable(
        "milestone1",
        month_width * 4 - milestone_size / 2, y_offset + activity_height / 2 - milestone_size / 2,
        milestone_size, milestone_size,
        "",
        shape_type='DIAMOND',
        fill_color=MockColor(255, 200, 0),
        line_color=activity_line,
        font_color=activity_font,
        font_size=9
    ))

    # QA activities
    y_offset2 = swimlane_top + swimlane_height + swimlane_gap + 20
    activities.append(MockPlotable(
        "act4",
        month_width * 2.5, y_offset2, month_width * 1.5, activity_height,
        "Unit Testing",
        shape_type='ROUNDED_RECTANGLE',
        fill_color=activity_orange,
        line_color=activity_line,
        font_color=activity_font,
        font_size=10
    ))

    activities.append(MockPlotable(
        "act5",
        month_width * 4.5, y_offset2, month_width * 2, activity_height,
        "Integration Testing",
        shape_type='ROUNDED_RECTANGLE',
        fill_color=activity_orange,
        line_color=activity_line,
        font_color=activity_font,
        font_size=10
    ))

    # Infrastructure activities
    y_offset3 = swimlane_top + 2 * (swimlane_height + swimlane_gap) + 20
    activities.append(MockPlotable(
        "act6",
        month_width * 1, y_offset3, month_width * 2, activity_height,
        "Setup CI/CD",
        shape_type='ROUNDED_RECTANGLE',
        fill_color=MockColor(180, 150, 200),
        line_color=activity_line,
        font_color=activity_font,
        font_size=10
    ))

    activities.append(MockPlotable(
        "act7",
        month_width * 5, y_offset3, month_width * 1.5, activity_height,
        "Deploy Prod",
        shape_type='ROUNDED_RECTANGLE',
        fill_color=MockColor(180, 150, 200),
        line_color=activity_line,
        font_color=activity_font,
        font_size=10
    ))

    return {
        "timelines": timelines,
        "swimlanes": swimlanes,
        "visual_activities": activities
    }


def render_to_powerpoint(plotables, output_path):
    """Render plotables to PowerPoint (simulating the actual renderer)"""

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Project Timeline - 2024"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(18)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Calculate dimensions and scale
    min_left = min_top = float('inf')
    max_right = max_bottom = float('-inf')

    for layer_plotables in plotables.values():
        for plotable in layer_plotables:
            min_left = min(min_left, plotable.left)
            min_top = min(min_top, plotable.top)
            max_right = max(max_right, plotable.get_right())
            max_bottom = max(max_bottom, plotable.get_bottom())

    visual_width = max_right - min_left
    visual_height = max_bottom - min_top

    # Calculate scale to fit on slide
    margin = Inches(0.5)
    title_space = Inches(0.7)
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    available_width = slide_width - 2 * margin
    available_height = slide_height - 2 * margin - title_space

    scale_x = available_width / visual_width if visual_width > 0 else 1
    scale_y = available_height / visual_height if visual_height > 0 else 1
    scale_factor = min(scale_x, scale_y)

    offset_left = margin
    offset_top = margin + title_space

    print(f"Visual dimensions: {visual_width:.1f} x {visual_height:.1f}")
    print(f"Scale factor: {scale_factor:.4f}")

    # Shape type mapping
    shape_map = {
        'RECTANGLE': MSO_SHAPE.RECTANGLE,
        'ROUNDED_RECTANGLE': MSO_SHAPE.ROUNDED_RECTANGLE,
        'DIAMOND': MSO_SHAPE.DIAMOND,
    }

    # Render each layer
    shape_count = 0
    for layer_name, layer_plotables in plotables.items():
        print(f"Rendering {layer_name}: {len(layer_plotables)} items")
        for plotable in layer_plotables:
            left = offset_left + plotable.left * scale_factor
            top = offset_top + plotable.top * scale_factor
            width = plotable.width * scale_factor
            height = plotable.height * scale_factor

            shape_type = shape_map.get(plotable.shape.value, MSO_SHAPE.RECTANGLE)
            shape = slide.shapes.add_shape(shape_type, left, top, width, height)

            # Set colors
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                plotable.format.fill_color.red,
                plotable.format.fill_color.green,
                plotable.format.fill_color.blue
            )

            line = shape.line
            line.color.rgb = RGBColor(
                plotable.format.line_color.red,
                plotable.format.line_color.green,
                plotable.format.line_color.blue
            )
            line.width = Pt(plotable.format.line_thickness)

            # Add text
            if plotable.text:
                text_frame = shape.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = plotable.text
                font = run.font
                font.size = Pt(plotable.format.font_size)
                font.color.rgb = RGBColor(
                    plotable.format.font_color.red,
                    plotable.format.font_color.green,
                    plotable.format.font_color.blue
                )

                p.alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            shape_count += 1

    print(f"\nTotal shapes created: {shape_count}")

    # Save
    prs.save(output_path)
    print(f"\n✓ Saved to: {output_path}")
    print(f"  File size: {os.path.getsize(output_path) / 1024:.1f} KB")

    return prs


if __name__ == "__main__":
    import os

    print("=" * 70)
    print("Creating Sample PowerPoint Visual")
    print("=" * 70)
    print()

    # Create visual data
    plotables = create_realistic_visual()

    # Render to PowerPoint
    output_path = "sample_project_timeline.pptx"
    render_to_powerpoint(plotables, output_path)

    print()
    print("=" * 70)
    print("✓ Complete! Open the file to view the visual:")
    print(f"  {os.path.abspath(output_path)}")
    print("=" * 70)
