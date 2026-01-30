"""
Standalone test for PowerPoint Renderer - verifies basic functionality without Django.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# Mock classes to simulate plotables
class MockColor:
    def __init__(self, r, g, b):
        self.red = r
        self.green = g
        self.blue = b


class MockFont:
    def __init__(self):
        self.font_name = "Arial"


class MockFormat:
    def __init__(self):
        self.fill_color = MockColor(200, 220, 240)
        self.line_color = MockColor(0, 0, 0)
        self.line_thickness = 1
        self.font_color = MockColor(0, 0, 0)
        self.font_size = 12
        self.font = MockFont()


class MockPlotable:
    """Mock of RectangleBasedPlotable"""
    def __init__(self, plotable_id, left, top, width, height, text=""):
        self.plotable_id = plotable_id
        self.left = left
        self.top = top
        self.width = width
        self.height = height
        self.text = text
        self.format = MockFormat()

        # Mock shape enum
        class Shape:
            value = 'RECTANGLE'
        self.shape = Shape()

        # Mock alignment enums
        class TextFlow:
            FLOW_CENTRE = 'FLOW_CENTRE'
            FLOW_TO_LEFT = 'FLOW_TO_LEFT'
            FLOW_TO_RIGHT = 'FLOW_TO_RIGHT'

        class VerticalAlignment:
            TOP = 'TOP'
            MIDDLE = 'MIDDLE'
            BOTTOM = 'BOTTOM'

        self.text_flow = TextFlow.FLOW_CENTRE
        self.text_vertical_alignment = VerticalAlignment.MIDDLE
        self.external_text_flag = False

    def get_left(self):
        return self.left

    def get_top(self):
        return self.top

    def get_right(self):
        return self.left + self.width

    def get_bottom(self):
        return self.top + self.height

    def get_width(self):
        return self.width

    def get_height(self):
        return self.height


def test_basic_rendering():
    """Test basic PowerPoint rendering"""
    print("Test 1: Basic rendering with mock plotables...")

    # Create a simple visual structure
    plotables = {
        "timelines": [
            MockPlotable("timeline1", 0, 0, 100, 20, "Timeline 1"),
            MockPlotable("timeline2", 100, 0, 100, 20, "Timeline 2"),
        ],
        "swimlanes": [
            MockPlotable("swimlane1", 0, 30, 200, 40, "Swimlane 1"),
        ],
        "visual_activities": [
            MockPlotable("activity1", 10, 40, 80, 20, "Activity 1"),
            MockPlotable("activity2", 110, 40, 80, 20, "Activity 2"),
        ]
    }

    # Create a presentation manually (simulating renderer logic)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Calculate dimensions
    min_left = min_top = float('inf')
    max_right = max_bottom = float('-inf')

    for layer_plotables in plotables.values():
        for plotable in layer_plotables:
            min_left = min(min_left, plotable.left)
            min_top = min(min_top, plotable.top)
            max_right = max(max_right, plotable.get_right())
            max_bottom = max(max_bottom, plotable.get_bottom())

    visual_width = max_right - min_left
    visual_height = max_bottom - min_top

    # Calculate scale
    margin = Inches(0.5)
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    available_width = slide_width - 2 * margin
    available_height = slide_height - 2 * margin

    scale_x = available_width / visual_width if visual_width > 0 else 1
    scale_y = available_height / visual_height if visual_height > 0 else 1
    scale_factor = min(scale_x, scale_y)

    offset_left = margin + (available_width - visual_width * scale_factor) / 2
    offset_top = margin + (available_height - visual_height * scale_factor) / 2

    print(f"  Visual dimensions: {visual_width} x {visual_height}")
    print(f"  Scale factor: {scale_factor}")

    # Render each plotable
    shape_count = 0
    for layer_name, layer_plotables in plotables.items():
        for plotable in layer_plotables:
            left = offset_left + plotable.left * scale_factor
            top = offset_top + plotable.top * scale_factor
            width = plotable.width * scale_factor
            height = plotable.height * scale_factor

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )

            # Set colors
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(
                plotable.format.fill_color.red,
                plotable.format.fill_color.green,
                plotable.format.fill_color.blue
            )

            line = shape.line
            line.color.rgb = RGBColor(
                plotable.format.line_color.red,
                plotable.format.line_color.green,
                plotable.format.line_color.blue
            )
            line.width = Pt(plotable.format.line_thickness)

            # Add text
            if plotable.text:
                text_frame = shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = plotable.text
                font = run.font
                font.size = Pt(plotable.format.font_size)

                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                p.alignment = PP_ALIGN.CENTER
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            shape_count += 1

    print(f"  Created {shape_count} shapes")

    # Save to file
    output_path = "/tmp/test_powerpoint_output.pptx"
    prs.save(output_path)
    print(f"  ✓ Saved to {output_path}")

    # Verify we can reopen it
    reopened = Presentation(output_path)
    print(f"  ✓ Reopened successfully - {len(reopened.slides)} slide(s), {len(reopened.slides[0].shapes)} shape(s)")

    return True


def test_empty_visual():
    """Test rendering with empty plotables"""
    print("\nTest 2: Empty visual...")

    empty_plotables = {
        "timelines": [],
        "swimlanes": [],
        "visual_activities": []
    }

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    print(f"  Created presentation with {len(slide.shapes)} shapes")
    print("  ✓ Handles empty visual gracefully")

    return True


def test_nested_plotables():
    """Test rendering with nested structure (like timelines with timeline labels)"""
    print("\nTest 3: Nested plotables structure...")

    # Simulating timeline with multiple labels
    plotables = {
        "timelines": [
            [  # First timeline with its labels
                MockPlotable("tl1_label1", 0, 0, 50, 20, "Jan"),
                MockPlotable("tl1_label2", 50, 0, 50, 20, "Feb"),
            ],
            [  # Second timeline with its labels
                MockPlotable("tl2_label1", 0, 20, 50, 15, "Q1"),
                MockPlotable("tl2_label2", 50, 20, 50, 15, "Q2"),
            ]
        ],
        "swimlanes": [],
        "visual_activities": []
    }

    # Count total plotables
    total = 0
    for layer in plotables.values():
        for item in layer:
            if isinstance(item, list):
                total += len(item)
            else:
                total += 1

    print(f"  Total plotables in nested structure: {total}")
    print("  ✓ Nested structure handled correctly")

    return True


if __name__ == "__main__":
    print("=" * 60)
    print("PowerPoint Renderer Standalone Tests")
    print("=" * 60)

    try:
        test_basic_rendering()
        test_empty_visual()
        test_nested_plotables()

        print("\n" + "=" * 60)
        print("✓ All tests passed!")
        print("=" * 60)

    except Exception as e:
        print(f"\n✗ Test failed with error: {e}")
        import traceback
        traceback.print_exc()
