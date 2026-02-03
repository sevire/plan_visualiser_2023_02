from abc import ABC, abstractmethod
from typing import Iterable, Dict, List, Optional, Tuple
from plan_visual_django.services.visual.rendering.plotables import RectangleBasedPlotable, Plotable
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE


class VisualRenderer(ABC):
    """
    Base class for all renderers that convert Plotables into a target medium.

    Architecture:
    - Takes plotables organized by layer (timelines, swimlanes, activities)
    - Layer order determines z-index/stacking order in the rendered output
    - Handles arbitrary nesting depth of plotables (e.g., timelines contain timeline labels)
    - Each subclass implements render_plotable() for medium-specific rendering
    """

    def render_from_iterable(self, visual_plotables: Dict[str, Iterable]):
        """
        Main entry point for rendering. Takes plotables organized by layer:
        {
            "timelines": [...],         # Background layer
            "swimlanes": [...],         # Middle layer
            "visual_activities": [...]  # Foreground layer
        }

        Layer order determines z-index in final output. Iterates through layers
        and recursively renders nested structures.

        :param visual_plotables: Dictionary mapping layer names to iterables of Plotables
        :return: Renderer-specific output (defined by _finalize_render())
        """
        for layer_name, plotable_iterable in visual_plotables.items():
            self._render_iterable(plotable_iterable)
        return self._finalize_render()

    def _render_iterable(self, plotable_iterable: Plotable | Iterable):
        """
        Recursively renders nested iterables of plotables.
        Handles arbitrary nesting depth (e.g., timeline -> timeline_labels).

        :param plotable_iterable: Either a single Plotable or an iterable containing
                                  Plotables and/or nested iterables
        """
        if isinstance(plotable_iterable, Plotable):
            self.render_plotable(plotable_iterable)
        elif isinstance(plotable_iterable, Iterable):
            for item in plotable_iterable:
                self._render_iterable(item)
        else:
            raise TypeError(f"Expected Plotable or Iterable, got {type(plotable_iterable)}")

    def _calculate_visual_bounds(self, visual_plotables: Dict[str, Iterable]):
        """
        Calculate the total dimensions of the visual by examining all plotables.
        
        :param visual_plotables: Dictionary of layer_name -> plotables
        :return: Tuple of (min_left, min_top, max_right, max_bottom, width, height)
        """
        min_left = float('inf')
        min_top = float('inf')
        max_right = float('-inf')
        max_bottom = float('-inf')

        # Updates visual bounds based on plotable dimensions
        def update_bounds(plotable: Plotable):
            nonlocal min_left, min_top, max_right, max_bottom
            min_left = min(min_left, plotable.get_left())
            min_top = min(min_top, plotable.get_top())
            max_right = max(max_right, plotable.get_right())
            max_bottom = max(max_bottom, plotable.get_bottom())

        def traverse(item):
            if isinstance(item, Plotable):
                update_bounds(item)
            elif isinstance(item, Iterable):
                for sub_item in item:
                    traverse(sub_item)

        for layer_plotables in visual_plotables.values():
            traverse(layer_plotables)

        width = max_right - min_left if max_right != float('-inf') else 0
        height = max_bottom - min_top if max_bottom != float('-inf') else 0

        return min_left, min_top, max_right, max_bottom, width, height

    def render_plotable(self, plotable: Plotable):
        """
        Render a single plotable to the target medium.
        Must be implemented by each subclass for their specific output format.

        :param plotable: The Plotable object to render
        """
        pass

    def _decompose_plotable(self, item: Plotable) -> Tuple[Optional[dict], Optional[dict]]:
        """
        Decomposes a Plotable into its shape and text components.
        This provides a common way for all renderers to handle text as a separate
        entity from the shape, allowing for better control over overflow and positioning.

        :param item: The Plotable to decompose
        :return: Tuple of (shape_component, text_component)
        """
        if isinstance(item, RectangleBasedPlotable):
            # Shape component
            shape_component = {
                'type': 'shape',
                'id': item.plotable_id,
                'shape_name': item.shape.value,
                'top': item.top,
                'left': item.left,
                'width': item.width,
                'height': item.height,
                'fill_color': item.format.fill_color,
                'line_color': item.format.line_color,
                'line_thickness': item.format.line_thickness,
            }

            # Text component
            text_component = None
            if item.text:
                text_x, text_align = item.get_text_x()
                text_y, text_baseline = item.get_text_y()
                text_component = {
                    'type': 'text',
                    'id': f"{item.plotable_id}-text",
                    'text': item.text,
                    'x': text_x,
                    'y': text_y,
                    'text_align': text_align,
                    'text_baseline': text_baseline,
                    'font_color': item.format.font_color,
                    'font_size': round(item.format.font_size),
                    'font_name': item.format.font.font_name,
                    'text_flow': item.text_flow,
                    'vertical_alignment': item.text_vertical_alignment,
                    # Original plotable dimensions for reference
                    'parent_top': item.top,
                    'parent_left': item.left,
                    'parent_width': item.width,
                    'parent_height': item.height,
                }
            return shape_component, text_component
        return None, None

    def _finalize_render(self):
        """
        Called after all plotables have been rendered.
        Override to return medium-specific output.

        :return: The final rendered output (format depends on subclass)
        """
        return None


class CanvasRenderer(VisualRenderer):
    """
    Renders plotables as JSON data structures for HTML Canvas rendering in the browser.

    This renderer doesn't perform actual drawing - it creates data structures that
    are serialized to JSON and sent to the browser, where JavaScript renders them
    onto HTML canvas elements.

    Output format: Dictionary mapping layer names to lists of canvas objects
    {
        "timelines": [{shape_data}, {shape_data}, ...],
        "swimlanes": [{shape_data}, {shape_data}, ...],
        "visual_activities": [{shape_data}, {shape_data}, ...]
    }
    """

    def __init__(self):
        """Initialize canvas renderer with empty output structure."""
        self.canvas_objects = {}
        self.current_layer = None

    def render_from_iterable(self, visual_plotables: Dict[str, Iterable]):
        """
        Override to track which layer we're rendering for proper organization.

        :param visual_plotables: Dictionary of layer_name -> plotables
        :return: Dictionary of layer_name -> rendered canvas objects
        """
        for layer_name, plotable_iterable in visual_plotables.items():
            self.current_layer = layer_name
            self.canvas_objects[layer_name] = []
            self._render_iterable(plotable_iterable)

        return self.canvas_objects

    def _render_iterable(self, plotable_iterable: Plotable | Iterable):
        """
        Override to accumulate rendered objects in current layer.
        """
        if isinstance(plotable_iterable, Plotable):
            rendered_objects = self.render_plotable(plotable_iterable)
            self.canvas_objects[self.current_layer].extend(rendered_objects)
        elif isinstance(plotable_iterable, Iterable):
            for item in plotable_iterable:
                self._render_iterable(item)
        else:
            raise TypeError(f"Expected Plotable or Iterable, got {type(plotable_iterable)}")

    def render_plotable(self, item: Plotable):
        """
        Converts a Plotable into JSON-serializable canvas rendering data.

        Note: A single Plotable may produce multiple canvas objects (e.g., separate
        objects for shape and text). Text positioning logic is calculated here since
        the browser only receives flat drawing instructions.

        :param item: The Plotable to render
        :return: List of canvas object dictionaries
        """
        shape_component, text_component = self._decompose_plotable(item)
        if not shape_component:
            return []

        canvas_rendered_objects = [
            {
                "plotable_id": shape_component['id'],
                'shape_type': 'rectangle',
                'shape_name': shape_component['shape_name'],
                'shape_plot_dims': {
                    'top': shape_component['top'],
                    'left': shape_component['left'],
                    'width': shape_component['width'],
                    'height': shape_component['height'],
                },
                'fill_color': shape_component['fill_color'].to_dict(),
                'stroke_color': shape_component['line_color'].to_dict(),
                'stoke_line_width': shape_component['line_thickness'],
            }
        ]

        if text_component:
            canvas_rendered_objects.append({
                "plotable_id": text_component['id'],
                'shape_type': 'text',
                'text': text_component['text'],
                'shape_name': None,
                'shape_plot_dims': {
                    'x': text_component['x'],
                    'y': text_component['y'],
                    'text_align': text_component['text_align'],
                    'text_baseline': text_component['text_baseline']
                },
                'fill_color': text_component['font_color'].to_dict(),
                'font_size': text_component['font_size']
            })

        return canvas_rendered_objects


class PowerPointRenderer(VisualRenderer):
    """
    Renders plotables as shapes on a PowerPoint slide.

    This renderer creates actual PowerPoint shapes using python-pptx library,
    drawing directly onto a slide with proper layering (timelines at back,
    activities at front).

    The coordinate system is converted from the visual's coordinate system
    (arbitrary units) to PowerPoint's coordinate system (EMUs/Inches).
    """

    def __init__(self, presentation: Optional[Presentation] = None, slide_index: int = None):
        """
        Initialize PowerPoint renderer.

        :param presentation: Optional existing Presentation. If None, creates new one.
        :param slide_index: Optional slide index to use. If None, adds a new blank slide.
        """
        if presentation is None:
            self.presentation = Presentation()
            # Add a blank slide layout (layout 6 is typically blank)
            blank_slide_layout = self.presentation.slide_layouts[6]
            self.slide = self.presentation.slides.add_slide(blank_slide_layout)
        else:
            self.presentation = presentation
            if slide_index is not None:
                self.slide = self.presentation.slides[slide_index]
            else:
                blank_slide_layout = self.presentation.slide_layouts[6]
                self.slide = self.presentation.slides.add_slide(blank_slide_layout)

        # Will be set during rendering to calculate coordinate conversion
        self.visual_width = None
        self.visual_height = None
        self.scale_factor = None

    def render_from_iterable(self, visual_plotables: Dict[str, Iterable]):
        """
        Override to calculate visual dimensions and scale factor before rendering.

        :param visual_plotables: Dictionary of layer_name -> plotables
        :return: The Presentation object with rendered slide
        """
        # Calculate the bounding box of all plotables to determine scale
        _, _, _, _, self.visual_width, self.visual_height = self._calculate_visual_bounds(visual_plotables)

        # Calculate scale factor to fit on slide with some margin
        margin = Inches(0.5)
        
        # Use presentation slide dimensions instead of hardcoded constants
        available_width = self.presentation.slide_width - 2 * margin
        available_height = self.presentation.slide_height - 2 * margin

        scale_x = available_width / self.visual_width if self.visual_width > 0 else 1
        scale_y = available_height / self.visual_height if self.visual_height > 0 else 1

        # Use the smaller scale to ensure it fits
        self.scale_factor = min(scale_x, scale_y)

        # Store offset to align to top with margin
        self.offset_left = margin + (available_width - self.visual_width * self.scale_factor) / 2
        self.offset_top = margin

        # Now render with the parent implementation
        super().render_from_iterable(visual_plotables)

        return self.presentation

    def _convert_to_pptx_coords(self, left: float, top: float, width: float, height: float):
        """
        Convert visual coordinates to PowerPoint coordinates.

        :return: (left, top, width, height) in PowerPoint units (Inches)
        """
        pptx_left = self.offset_left + left * self.scale_factor
        pptx_top = self.offset_top + top * self.scale_factor
        pptx_width = width * self.scale_factor
        pptx_height = height * self.scale_factor

        return pptx_left, pptx_top, pptx_width, pptx_height

    def _get_shape_type(self, shape_name: str):
        """
        Map plotable shape names to PowerPoint shape types.

        :param shape_name: Name from PlotableShapeName enum
        :return: MSO_SHAPE constant
        """
        shape_mapping = {
            'RECTANGLE': MSO_SHAPE.RECTANGLE,
            'ROUNDED_RECTANGLE': MSO_SHAPE.ROUNDED_RECTANGLE,
            'DIAMOND': MSO_SHAPE.DIAMOND,
            'ISOSCELES_TRIANGLE': MSO_SHAPE.ISOSCELES_TRIANGLE,
            'BULLET': MSO_SHAPE.OVAL,  # Use oval for bullets
        }
        return shape_mapping.get(shape_name, MSO_SHAPE.RECTANGLE)

    def render_plotable(self, item: Plotable):
        """
        Render a single plotable as a PowerPoint shape.

        :param item: The Plotable to render
        """
        shape_component, text_component = self._decompose_plotable(item)
        if not shape_component:
            return

        # 1. Render the shape
        left, top, width, height = self._convert_to_pptx_coords(
            shape_component['left'], shape_component['top'],
            shape_component['width'], shape_component['height']
        )

        shape_type = self._get_shape_type(shape_component['shape_name'])
        shape = self.slide.shapes.add_shape(
            shape_type,
            left, top, width, height
        )

        # Set fill color
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(
            shape_component['fill_color'].red,
            shape_component['fill_color'].green,
            shape_component['fill_color'].blue
        )

        # Set line color and thickness
        line = shape.line
        line.color.rgb = RGBColor(
            shape_component['line_color'].red,
            shape_component['line_color'].green,
            shape_component['line_color'].blue
        )
        line.width = Pt(shape_component['line_thickness'])

            # 2. Render the text as a separate textbox if present
        if text_component:
            # We want to support overflow, so we'll make the textbox position and width
            # based on the parent shape and text flow.

            # We use the parent's dimensions as a baseline but will adjust for overflow.
            t_left, t_top, t_width, t_height = self._convert_to_pptx_coords(
                text_component['parent_left'], text_component['parent_top'],
                text_component['parent_width'], text_component['parent_height']
            )

            # To support overflow, we calculate a "generous" width for the textbox
            # if it's set to flow beyond the shape.
            from plan_visual_django.models import VisualActivity

            # Calculate slide boundaries in EMUs
            slide_width_emu = self.presentation.slide_width
            
            if text_component['text_flow'] == VisualActivity.TextFlow.FLOW_TO_RIGHT:
                # Keep t_left as start of shape, but allow it to be as wide as possible
                # until the right edge of the slide.
                t_width = slide_width_emu - t_left
            elif text_component['text_flow'] == VisualActivity.TextFlow.FLOW_TO_LEFT:
                # End at the right edge of the shape, but allow it to be as wide
                # as possible until the left edge of the slide.
                original_right = t_left + t_width
                t_left = 0
                t_width = original_right
            elif text_component['text_flow'] == VisualActivity.TextFlow.FLOW_CENTRE:
                # Center on the shape, allow overflow in both directions up to slide edges.
                # To keep it perfectly centered over the shape, we make the textbox
                # symmetric around the shape's center, constrained by the closest slide edge.
                centre_x = t_left + t_width / 2
                
                # Calculate distance to both edges
                dist_to_left = centre_x
                dist_to_right = slide_width_emu - centre_x
                
                # Use the smaller distance to maintain symmetry
                max_half_width = min(dist_to_left, dist_to_right)
                
                t_left = centre_x - max_half_width
                t_width = 2 * max_half_width
            else:
                # For clipped or within shape, we don't change t_left/t_width from shape dimensions
                pass

            # Ensure width is at least something small to avoid errors
            t_width = max(t_width, Inches(0.1))

            # Create a textbox
            textbox = self.slide.shapes.add_textbox(t_left, t_top, t_width, t_height)
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # Set word wrap to False to allow overflow beyond the width
            text_frame.word_wrap = False
            # Allow the shape to grow to fit text if needed (though we might want to control it)
            # text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text_component['text']

            # Set font properties
            font = run.font
            font.size = Pt(text_component['font_size'])
            font.color.rgb = RGBColor(
                text_component['font_color'].red,
                text_component['font_color'].green,
                text_component['font_color'].blue
            )
            if text_component['font_name']:
                font.name = text_component['font_name']

            # Set text alignment based on text_flow
            from plan_visual_django.models import VisualActivity
            if text_component['text_flow'] == VisualActivity.TextFlow.FLOW_CENTRE:
                p.alignment = PP_ALIGN.CENTER
            elif text_component['text_flow'] == VisualActivity.TextFlow.FLOW_TO_LEFT:
                p.alignment = PP_ALIGN.RIGHT
            else:  # FLOW_TO_RIGHT or default
                p.alignment = PP_ALIGN.LEFT

            # Set vertical alignment
            if text_component['vertical_alignment'] == VisualActivity.VerticalAlignment.MIDDLE:
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            elif text_component['vertical_alignment'] == VisualActivity.VerticalAlignment.BOTTOM:
                text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
            else:  # TOP or default
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
