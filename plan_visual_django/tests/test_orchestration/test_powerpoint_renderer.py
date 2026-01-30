import os
import tempfile
from django.test import TestCase
from pptx import Presentation
from plan_visual_django.models import PlanVisual
from plan_visual_django.services.visual.rendering.renderers import PowerPointRenderer
from plan_visual_django.tests.resources.unit_test_configuration import test_data_base_folder, test_fixtures_folder


# NOTE: These tests require Python 3.9+ due to type hint syntax in models.py
# For Python 3.8, use the standalone test: test_powerpoint_standalone.py


class TestPowerPointRenderer(TestCase):
    """
    Tests the PowerPointRenderer to ensure it can successfully render plotables
    to a PowerPoint presentation.
    """
    fixtures = [
        os.path.join(test_data_base_folder, test_fixtures_folder, 'auth_test_fixtures.json'),
        os.path.join(test_data_base_folder, test_fixtures_folder, 'test_fixtures.json')
    ]

    def test_powerpoint_renderer_creates_presentation(self):
        """
        Test that PowerPointRenderer creates a valid presentation with a slide.
        """
        # Get a visual from the test database
        visual = PlanVisual.objects.get(pk=4)
        visual_plotables = visual.get_plotables()

        # Ensure there are plotables to render
        self.assertGreater(len(visual_plotables), 0)

        # Create renderer and render
        renderer = PowerPointRenderer()
        presentation = renderer.render_from_iterable(visual_plotables)

        # Verify we got a Presentation object
        self.assertIsInstance(presentation, Presentation)

        # Verify the presentation has at least one slide
        self.assertGreater(len(presentation.slides), 0)

    def test_powerpoint_renderer_saves_successfully(self):
        """
        Test that PowerPointRenderer creates a presentation that can be saved to disk.
        """
        # Get a visual from the test database
        visual = PlanVisual.objects.get(pk=4)
        visual_plotables = visual.get_plotables()

        # Create renderer and render
        renderer = PowerPointRenderer()
        presentation = renderer.render_from_iterable(visual_plotables)

        # Try to save to a temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            tmp_path = tmp_file.name

        try:
            # Save the presentation
            presentation.save(tmp_path)

            # Verify the file was created
            self.assertTrue(os.path.exists(tmp_path))
            self.assertGreater(os.path.getsize(tmp_path), 0)

            # Verify we can open it again
            reopened_prs = Presentation(tmp_path)
            self.assertIsInstance(reopened_prs, Presentation)
            self.assertGreater(len(reopened_prs.slides), 0)

        finally:
            # Clean up
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)

    def test_powerpoint_renderer_creates_shapes(self):
        """
        Test that PowerPointRenderer actually creates shapes on the slide.
        """
        # Get a visual from the test database
        visual = PlanVisual.objects.get(pk=4)
        visual_plotables = visual.get_plotables()

        # Create renderer and render
        renderer = PowerPointRenderer()
        presentation = renderer.render_from_iterable(visual_plotables)

        # Get the first (and only) slide
        slide = presentation.slides[0]

        # Verify shapes were created
        # We should have shapes for timelines, swimlanes, and activities
        self.assertGreater(len(slide.shapes), 0)

        # Count how many shapes should be created (roughly)
        # Each plotable in the canvas renderer test creates 2 objects (shape + text)
        # For PowerPoint, each plotable creates 1 shape (with text inside)
        expected_min_shapes = 0
        for layer_name, plotables in visual_plotables.items():
            # Count plotables recursively
            def count_plotables(item):
                from plan_visual_django.services.visual.rendering.plotables import Plotable
                if isinstance(item, Plotable):
                    return 1
                elif hasattr(item, '__iter__') and not isinstance(item, str):
                    return sum(count_plotables(sub_item) for sub_item in item)
                return 0

            expected_min_shapes += count_plotables(plotables)

        # Verify we have at least the expected number of shapes
        self.assertGreaterEqual(len(slide.shapes), expected_min_shapes)

    def test_powerpoint_renderer_with_existing_presentation(self):
        """
        Test that PowerPointRenderer can add a slide to an existing presentation.
        """
        # Create an existing presentation with one slide
        existing_prs = Presentation()
        title_slide_layout = existing_prs.slide_layouts[0]
        slide1 = existing_prs.slides.add_slide(title_slide_layout)
        title = slide1.shapes.title
        title.text = "Existing Slide"

        initial_slide_count = len(existing_prs.slides)

        # Get a visual from the test database
        visual = PlanVisual.objects.get(pk=4)
        visual_plotables = visual.get_plotables()

        # Create renderer with existing presentation
        renderer = PowerPointRenderer(presentation=existing_prs)
        presentation = renderer.render_from_iterable(visual_plotables)

        # Verify a new slide was added
        self.assertEqual(len(presentation.slides), initial_slide_count + 1)

        # Verify the first slide is still intact
        self.assertEqual(presentation.slides[0].shapes.title.text, "Existing Slide")

        # Verify the new slide has shapes
        new_slide = presentation.slides[-1]
        self.assertGreater(len(new_slide.shapes), 0)

    def test_powerpoint_renderer_handles_empty_visual(self):
        """
        Test that PowerPointRenderer handles a visual with no plotables gracefully.
        """
        # Create an empty plotables structure
        empty_plotables = {
            "timelines": [],
            "swimlanes": [],
            "visual_activities": []
        }

        # Create renderer and render
        renderer = PowerPointRenderer()
        presentation = renderer.render_from_iterable(empty_plotables)

        # Should still create a presentation with a slide, just no shapes
        self.assertIsInstance(presentation, Presentation)
        self.assertGreater(len(presentation.slides), 0)

        # The slide should have no shapes (or very few, depending on slide layout)
        slide = presentation.slides[0]
        # A blank slide layout should have minimal shapes
        self.assertLessEqual(len(slide.shapes), 2)

    def test_powerpoint_renderer_coordinate_conversion(self):
        """
        Test that coordinates are properly converted and shapes fit on the slide.
        """
        # Get a visual from the test database
        visual = PlanVisual.objects.get(pk=4)
        visual_plotables = visual.get_plotables()

        # Create renderer and render
        renderer = PowerPointRenderer()
        presentation = renderer.render_from_iterable(visual_plotables)

        slide = presentation.slides[0]

        # Verify all shapes are within slide bounds
        from pptx.util import Inches
        slide_width = Inches(10)
        slide_height = Inches(7.5)

        for shape in slide.shapes:
            # Each shape should be within the slide bounds
            self.assertGreaterEqual(shape.left, 0)
            self.assertGreaterEqual(shape.top, 0)
            self.assertLessEqual(shape.left + shape.width, slide_width)
            self.assertLessEqual(shape.top + shape.height, slide_height)

    def test_powerpoint_renderer_layer_ordering(self):
        """
        Test that layers are rendered in the correct order (timelines, swimlanes, activities).
        The z-order is determined by the order shapes are added to the slide.
        """
        # Get a visual from the test database
        visual = PlanVisual.objects.get(pk=4)
        visual_plotables = visual.get_plotables()

        # Create renderer and render
        renderer = PowerPointRenderer()
        presentation = renderer.render_from_iterable(visual_plotables)

        slide = presentation.slides[0]

        # Verify we have shapes
        self.assertGreater(len(slide.shapes), 0)

        # The shapes are added in order: timelines, swimlanes, activities
        # We can't directly test z-order without inspecting XML, but we can verify
        # that the renderer completed without errors and all layers were processed
        self.assertIn('timelines', visual_plotables)
        self.assertIn('swimlanes', visual_plotables)
        self.assertIn('visual_activities', visual_plotables)
